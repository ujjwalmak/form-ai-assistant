#!/usr/bin/env python3
"""FormAssist — Reflexion deck (Aurora-Glass), 12 slides, 16:9, mit Sprechertext."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls

import os
_HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(_HERE, "assets")
OUT = os.path.join(_HERE, "..", "FormAssist_Reflexion.pptx")

EMU_PX = 6350  # 12192000 / 1920
def PX(v): return Emu(int(v * EMU_PX))

# ---- Design-Tokens (aus reflexion.html) ----
GRAD = [(0, "8B5CF6"), (48, "D946EF"), (100, "F472B6")]
TITLE_GRAD = [(10, "FFFFFF"), (42, "E9D5FF"), (72, "F0ABFC"), (100, "F9A8D4")]
FG, DIM, LILAC, PINKISH = "ECEBF2", "B0AAC2", "E9D5FF", "F0ABFC"
SG, SGM, SGS = "Space Grotesk", "Space Grotesk Medium", "Space Grotesk SemiBold"
IN, INM, INS = "Inter", "Inter Medium", "Inter SemiBold"
ANG_110, ANG_V = 1200000, 5400000

FILLS = ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill")

def _gradfill_xml(stops, ang):
    gs = "".join(f'<a:gs pos="{int(p*1000)}"><a:srgbClr val="{c}"/></a:gs>' for p, c in stops)
    return (f'<a:gradFill {nsdecls("a")} rotWithShape="1"><a:gsLst>{gs}</a:gsLst>'
            f'<a:lin ang="{ang}" scaled="1"/></a:gradFill>')

def _set_sp_fill(shape, xml):
    spPr = shape._element.spPr
    for tag in FILLS:
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    ln = spPr.find(qn("a:ln"))
    el = parse_xml(xml)
    if ln is not None:
        ln.addprevious(el)
    else:
        spPr.append(el)

def solid_alpha(shape, hexc, alpha_pct):
    _set_sp_fill(shape, f'<a:solidFill {nsdecls("a")}><a:srgbClr val="{hexc}">'
                        f'<a:alpha val="{int(alpha_pct*1000)}"/></a:srgbClr></a:solidFill>')

def grad_fill(shape, stops=GRAD, ang=ANG_110):
    _set_sp_fill(shape, _gradfill_xml(stops, ang))

def _get_ln(shape):
    spPr = shape._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = parse_xml(f'<a:ln {nsdecls("a")}/>')
        spPr.append(ln)
    for tag in FILLS:
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    return ln

def line_alpha(shape, hexc, alpha_pct, w_pt):
    ln = _get_ln(shape)
    ln.set("w", str(int(w_pt * 12700)))
    ln.insert(0, parse_xml(f'<a:solidFill {nsdecls("a")}><a:srgbClr val="{hexc}">'
                           f'<a:alpha val="{int(alpha_pct*1000)}"/></a:srgbClr></a:solidFill>'))

def grad_line(shape, stops=GRAD, ang=ANG_110, w_pt=1.2):
    ln = _get_ln(shape)
    ln.set("w", str(int(w_pt * 12700)))
    ln.insert(0, parse_xml(_gradfill_xml(stops, ang)))

def no_line(shape):
    _get_ln(shape).append(parse_xml(f'<a:noFill {nsdecls("a")}/>'))

def glow(shape, hexc="D946EF", alpha_pct=55, rad_pt=6):
    shape._element.spPr.append(parse_xml(
        f'<a:effectLst {nsdecls("a")}><a:glow rad="{int(rad_pt*12700)}">'
        f'<a:srgbClr val="{hexc}"><a:alpha val="{int(alpha_pct*1000)}"/></a:srgbClr>'
        f'</a:glow></a:effectLst>'))

def no_shadow(shape):
    shape.shadow.inherit = False

def grad_run(run, stops=GRAD, ang=ANG_110):
    rPr = run._r.get_or_add_rPr()
    for tag in FILLS:
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    el = parse_xml(_gradfill_xml(stops, ang))
    ln = rPr.find(qn("a:ln"))
    if ln is not None:
        ln.addnext(el)
    else:
        rPr.insert(0, el)

def add_runs(p, segs, size, default_font=IN, default_color=FG):
    for s in segs:
        r = p.add_run()
        r.text = s["t"]
        r.font.size = Pt(s.get("sz", size))
        r.font.name = s.get("f", default_font)
        r.font.bold = s.get("b", False)
        if s.get("spc"):
            r._r.get_or_add_rPr().set("spc", str(int(s["spc"])))
        if s.get("grad"):
            grad_run(r, s["grad"] if isinstance(s["grad"], list) else GRAD)
        else:
            r.font.color.rgb = RGBColor.from_string(s.get("c", default_color))
    return p

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(PX(x), PX(y), PX(w), PX(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb, tf

def text(slide, x, y, w, h, segs, size, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=None,
         default_font=IN, default_color=FG):
    tb, tf = textbox(slide, x, y, w, h, anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    if ls:
        p.line_spacing = ls
    add_runs(p, segs, size, default_font, default_color)
    return tb

def est_lines(segs, size, inner_w):
    n = sum(len(s["t"]) for s in segs)
    cpl = inner_w / (size * 2 * 0.48)
    return 1 if n <= cpl else 2

def chip_h(segs, size, w):
    lines = est_lines(segs, size, w - 56 - 26)
    return (58 if size >= 12 else 54) if lines == 1 else (88 if size >= 12 else 84)

def chip(slide, x, y, w, segs, size=12.5, hi=False, dot=True, ls=1.3, align=PP_ALIGN.LEFT):
    h = chip_h(segs, size, w)
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sh.adjustments[0] = min(0.5, 14.0 / h)
    solid_alpha(sh, "FFFFFF", 6.0 if hi else 3.5)
    if hi:
        grad_line(sh, w_pt=1.4)
    else:
        line_alpha(sh, "FFFFFF", 7, 1.0)
    no_shadow(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left, tf.margin_right = PX(56 if dot else 26), PX(26)
    tf.margin_top = tf.margin_bottom = PX(6)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = ls
    add_runs(p, segs, size)
    if dot:
        d = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x + 24), PX(y + h/2 - 5.5), PX(11), PX(11))
        d.adjustments[0] = 0.33
        grad_fill(d)
        no_line(d)
        glow(d)
        no_shadow(d)
    return y + h

def bullet_stack(slide, x, w, specs, top, bottom=945, gap=14, bias=0.40, cap=120):
    """specs: list of (segs, hi) — vertikal ausbalanciert im Kartenbereich."""
    total = sum(chip_h(segs, 12.5, w) for segs, _ in specs) + gap * (len(specs) - 1)
    y = top + max(0, min(cap, (bottom - top - total) * bias))
    for segs, hi in specs:
        y = chip(slide, x, y, w, segs, hi=hi) + gap
    return y - gap

def eyebrow(slide, x, y, txt, color=LILAC, size=10, align=PP_ALIGN.LEFT, w=900):
    return text(slide, x, y, w, 30, [{"t": txt.upper(), "f": INS, "c": color, "spc": 110}],
                size, align=align)

def stat_tile(slide, x, y, w, h, num, label):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sh.adjustments[0] = min(0.5, 18.0 / h)
    solid_alpha(sh, "FFFFFF", 3.5)
    line_alpha(sh, "FFFFFF", 8, 1.0)
    no_shadow(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = PX(12)
    tf.margin_top = tf.margin_bottom = PX(8)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    add_runs(p1, [{"t": num, "f": SG, "b": True, "grad": True}], 30)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    add_runs(p2, [{"t": label, "f": IN, "c": DIM}], 10.5)

def mini_card(slide, x, y, w, h, title, sub, hi=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sh.adjustments[0] = min(0.5, 16.0 / h)
    solid_alpha(sh, "FFFFFF", 6.0 if hi else 3.5)
    if hi:
        grad_line(sh, w_pt=1.4)
    else:
        line_alpha(sh, "FFFFFF", 8, 1.0)
    no_shadow(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = PX(10)
    tf.margin_top = tf.margin_bottom = PX(6)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    add_runs(p1, [{"t": title, "f": INS, "c": "FFFFFF"}], 13)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(3)
    add_runs(p2, [{"t": sub, "f": IN, "c": DIM}], 9.5)

def hud(slide, idx, total=12):
    if idx > 0:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PX(0), PX(0), PX(1920 * idx / (total - 1)), PX(6))
        grad_fill(bar)
        no_line(bar)
        no_shadow(bar)
    text(slide, 1650, 1032, 230, 30,
         [{"t": f"{idx+1:02d}", "f": INS, "c": "FFFFFF"}, {"t": f" / {total:02d}", "f": IN, "c": DIM}],
         9.5, align=PP_ALIGN.RIGHT)

def q_header(slide, num, question, key):
    text(slide, 220, 112, 400, 80, [{"t": num, "f": SG, "b": True, "grad": True}], 40)
    text(slide, 220, 196, 1480, 60, [{"t": question, "f": SGS, "c": "FFFFFF"}], 26)
    lines = est_lines([{"t": key}], 15.5, 1420)
    bar_h = 46 if lines == 1 else 72
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(221), PX(276), PX(5), PX(bar_h))
    bar.adjustments[0] = 0.5
    grad_fill(bar, ang=ANG_V)
    no_line(bar)
    no_shadow(bar)
    text(slide, 246, 278, 1434, 80, [{"t": key, "f": INS, "c": "F6F3FB"}], 15.5, ls=1.32)
    return 396 if lines == 1 else 410

def link_pill(slide, x, y, w, h, size=14):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    pill.adjustments[0] = 0.5
    grad_fill(pill)
    no_line(pill)
    no_shadow(pill)
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = PX(10)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_runs(p, [{"t": "🌐  ujjwalmak.github.io/form-ai-assistant", "f": INS, "c": "FFFFFF"}], size)

W = lambda t: {"t": t, "f": INS, "c": "FFFFFF"}   # weißes Semibold-Segment
N = lambda t: {"t": t}                              # normales Segment

# =====================================================================
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
prs.core_properties.title = "FormAssist — Reflexion · AI-Prototyping SS 2026"
prs.core_properties.author = "Maximilian Plitzko · Ujjwal Makkar"
BLANK = prs.slide_layouts[6]
NOTES = []

def new_slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(f"{ASSETS}/{bg}", 0, 0, Emu(12192000), Emu(6858000))
    return s

# ---------------------------------------------------------------- 1 · Titel
s = new_slide("bg_title.png")
s.shapes.add_picture(f"{ASSETS}/logo.png", PX(960 - 92), PX(148), PX(184), PX(184))
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(720), PX(352), PX(480), PX(46))
badge.adjustments[0] = 0.5
solid_alpha(badge, "D946EF", 12)
line_alpha(badge, "D946EF", 40, 1.0)
no_shadow(badge)
btf = badge.text_frame
btf.vertical_anchor = MSO_ANCHOR.MIDDLE
btf.margin_left = btf.margin_right = PX(10)
btf.margin_top = btf.margin_bottom = 0
bp = btf.paragraphs[0]
bp.alignment = PP_ALIGN.CENTER
add_runs(bp, [{"t": "REFLEXION · AI-PROTOTYPING SS 2026", "f": INS, "c": LILAC, "spc": 90}], 10.5)
text(s, 160, 418, 1600, 200, [{"t": "FormAssist", "f": SG, "b": True, "grad": TITLE_GRAD}],
     66, align=PP_ALIGN.CENTER)
text(s, 460, 634, 1000, 80,
     [{"t": "KI-Assistent, der komplexe Browser-Formulare versteht, erklärt und ausfüllt — datensparsam und nie ohne Freigabe.",
       "f": IN, "c": DIM}], 14.5, align=PP_ALIGN.CENTER, ls=1.4)
text(s, 360, 758, 1200, 40,
     [{"t": "Team  ", "f": INS, "c": "FFFFFF"},
      {"t": "Maximilian Plitzko · Ujjwal Makkar", "f": INS, "grad": True}],
     14, align=PP_ALIGN.CENTER)
text(s, 360, 806, 1200, 34,
     [{"t": "Prof. Dr. Sebastian Dünnebeil · FK07 · 02.07.2026", "f": IN, "c": DIM}],
     11.5, align=PP_ALIGN.CENTER)
hud(s, 0)
NOTES.append("""⏱ ~45 Sek · Sprecher: Maximilian (Einstieg)

Guten Tag zusammen! Wir sind Maximilian Plitzko und Ujjwal Makkar, und das ist unsere Reflexion zu FormAssist — unserem KI-Assistenten, der komplexe Browser-Formulare versteht, in einfacher Sprache erklärt und ausfüllt.

Zwei Dinge waren uns von Anfang an wichtig: Wir arbeiten datensparsam, und es wird niemals etwas ohne die ausdrückliche Freigabe des Nutzers abgesendet.

Wir gehen heute entlang der acht Leitfragen aus der Vorlesung — wir haben sie untereinander aufgeteilt. Zum Einstieg ein kurzer Blick auf die Eckdaten.""")

# ---------------------------------------------------------------- 2 · Eckdaten
s = new_slide("bg_content.png")
eyebrow(s, 220, 118, "Auf einen Blick")
text(s, 220, 152, 1200, 60, [{"t": "Eckdaten des Projekts", "f": SGS, "c": "FFFFFF"}], 26)
tiles = [("11", "Wochen · 09.04.–25.06.2026"), ("52", "Commits · 2er-Team"),
         ("7", "Module · Shadow-DOM-UI"), ("69", "Unit-Tests · Vitest"),
         ("~77 %", "Branch-Coverage + CI"), ("0", "Auto-Submits — harte Guardrail")]
for i, (num, lab) in enumerate(tiles):
    stat_tile(s, 220 + (i % 3) * 505, 240 + (i // 3) * 200, 470, 170, num, lab)
bullet_stack(s, 220, 1480, [
    ([W("Chrome-Extension (Manifest V3) — "), N("Vanilla JS, kein Build-Step, direkt als entpackte Erweiterung ladbar.")], False),
    ([W("Stack: "), N("Groq + OpenRouter (LLM mit Auto-Fallback) · Supabase · Vitest · GitHub Actions.")], False),
    ([W("Extras: "), N("autonomer Doku-Agent (Flask + Git-Hook) · Stakeholder-Webseite (MkDocs) — live.")], False),
], top=688, bottom=940, bias=0.5, cap=40)
hud(s, 1)
NOTES.append("""⏱ ~1 Min · Sprecher: Maximilian

Elf Wochen Projektlaufzeit, 52 Commits, zwei Personen — beide durchgängig aktiv.

Das Produkt ist eine Chrome-Extension nach Manifest V3, bewusst in Vanilla JS ohne Build-Step gebaut, in sieben Module gegliedert; die komplette UI läuft isoliert im Shadow DOM.

Als LLM-Provider nutzen wir Groq, mit OpenRouter als automatischem Fallback, dazu Supabase für den optionalen Sync. Qualität: 69 Unit-Tests mit rund 77 Prozent Branch-Coverage, die bei jedem Push in der CI laufen.

Zwei Extras, auf die wir noch zu sprechen kommen: ein eigener autonomer Doku-Agent, der per Git-Hook nach jedem Commit dokumentiert, und eine Stakeholder-Webseite, live auf GitHub Pages.

Die Null unten rechts ist übrigens Absicht: null automatische Submits — das ist unsere härteste Guardrail. Damit zur ersten Frage.""")

# ---------------------------------------------------------------- 3 · Q1
s = new_slide("bg_content.png")
top = q_header(s, "01", "Welches Problem wollten Sie lösen?",
               "Komplexe Formulare fehlerfrei ausfüllen — gebündelt, erklärt, datensparsam und nie automatisch abgesendet.")
bullet_stack(s, 220, 1480, [
    ([W("Behördendeutsch & Rechtssprache"), N(" führen zu Fehlern — komplexe Browser-Formulare sind schwer verständlich.")], False),
    ([W("Nötige Angaben liegen verstreut: "), N("Rentenversicherungsnummer auf Papier, Ausweis in der Cloud, Bescheinigungen in Primuss & E-Mail-Postfach.")], False),
    ([W("Unsere Lösung: "), N("alle Daten in "), W("einer"), N(" Anwendung bündeln und beim Ausfüllen geführt unterstützen.")], True),
    ([W("Kern ist die korrekte Eingabe: "), N("Felder semantisch lesen · einfach erklären · datensparsam übertragen · "), W("nie Auto-Submit.")], False),
], top=top)
hud(s, 2)
NOTES.append("""⏱ ~1,5 Min · Sprecher: Maximilian

Ausgangspunkt war ein Problem, das jeder kennt: komplexe Browser-Formulare — Anträge, Immatrikulation, Behördenportale. Behördendeutsch und Rechtssprache führen zu Fehlern.

Dazu kommt: Die nötigen Angaben liegen verstreut. Die Rentenversicherungsnummer steht auf Papier, der Ausweis liegt in der Cloud, die Immatrikulationsbescheinigung in Primuss, die Krankenkassen-Bestätigung im E-Mail-Postfach.

Unsere Lösung bündelt alle Daten in einer Anwendung und unterstützt beim Ausfüllen geführt.

Wichtig ist uns die Abgrenzung: Es geht nicht nur um Bequemlichkeit, sondern um die korrekte Eingabe. Felder werden semantisch ausgelesen — nicht bloß per HTML-Scan —, in einfacher Sprache erklärt und datensparsam übertragen. Und abgesendet wird grundsätzlich nie automatisch; das ist eine harte Guardrail im Code.""")

# ---------------------------------------------------------------- 4 · Q2
s = new_slide("bg_content.png")
top = q_header(s, "02", "Was war die größte Herausforderung?",
               "Formulare über sehr heterogene Seiten zuverlässig lesen und füllen — und Gelerntes wiederverwenden statt neu zu fragen.")
bullet_stack(s, 220, 1480, [
    ([W("Technisch: "), N("Shadow DOM, Custom-Widgets (Kendo), Datepicker-Libraries, cross-origin iFrames.")], False),
    ([W("Gelerntes nutzen statt neu fragen: "), N("mehrere Umbauten — Konfidenz-Schwelle → Field-by-Field → Batch-Prompt.")], False),
    ([W("Richtung & Mehrwert: "), N("Was bauen wir, wohin geht es — und wie entsteht echter Mehrwert?")], False),
], top=top, bottom=650, cap=20)
eyebrow(s, 220, 682, "War Story — die Provider-Odyssee")
fx, fw, fh, gap = 220, 331, 112, 52
subs = [("Anthropic", "Credit-Limits erreicht", False), ("Gemini", "Zwischenstopp", False),
        ("Groq", "Primär — schnell & stabil", True), ("OpenRouter", "Auto-Fallback bei 429/5xx", True)]
for i, (t1, t2, hi) in enumerate(subs):
    mini_card(s, fx + i * (fw + gap), 724, fw, fh, t1, t2, hi=hi)
    if i < 3:
        text(s, fx + i * (fw + gap) + fw, 724, gap, fh,
             [{"t": "→", "f": SG, "b": True, "grad": True}], 20,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
hud(s, 3)
NOTES.append("""⏱ ~2 Min · Sprecher: Maximilian, War Story: Ujjwal

Die größte Herausforderung hatte drei Ebenen.

Erstens die Technik: Formulare über völlig unterschiedliche Seiten zuverlässig lesen und füllen — Shadow DOM, Custom-Widgets wie Kendo, diverse Datepicker-Libraries, cross-origin iFrames.

Zweitens: den Agenten dazu bringen, Gelerntes wiederzuverwenden statt neu zu fragen. Das hat mehrere Umbauten gekostet — von der Konfidenz-Schwelle über einen Field-by-Field-Agenten bis zum heutigen Batch-Prompt.

Drittens die Richtung: Was bauen wir überhaupt, und wie entsteht echter Mehrwert?

[Übergabe an Ujjwal:] Die Provider-Odyssee habe ich selbst durchlaufen. Zuerst Anthropic angebunden — aber in Credit-Limits gelaufen. Dann Gemini, schließlich Groq. Jeder Wechsel hieß: Prompts, Modell-IDs und max_tokens neu justieren. Am Ende habe ich OpenRouter als automatischen Fallback eingebaut: Ein 429 oder 5xx bei Groq blockiert den Nutzer nicht mehr — der Wechsel passiert transparent im Hintergrund, mit einem kurzen Toast als Hinweis.""")

# ---------------------------------------------------------------- 5 · Q3
s = new_slide("bg_content.png")
top = q_header(s, "03", "Welche Rollenverteilung hatten Sie im Team?",
               "Zusammenarbeit auf Augenhöhe, keine festen Silos — Aufgaben wöchentlich nach Bedarf neu verteilt.")
bullet_stack(s, 220, 1480, [
    ([W("Keine starre Rollenfestlegung — "), N("wöchentliche Reviews mit Neuverteilung der Aufgaben.")], False),
    ([W("Das Commit-Bild bestätigt es: "), N("beide durchgängig aktiv (35 · 20 Commits) — geteilte Verantwortung statt Silos.")], False),
    ([W("Jeder arbeitet auch im Code des anderen — "), N("leichtere Reviews, kein Inselwissen.")], False),
    ([W("Ehrliche Einordnung: "), N("bei 2 Personen empfehlenswert — in größeren Teams kritisch zu hinterfragen.")], False),
], top=top)
hud(s, 4)
NOTES.append("""⏱ ~1,5 Min · Sprecher: Ujjwal

Wir haben bewusst auf Augenhöhe gearbeitet, ohne feste Silos. Statt starrer Rollen gab es wöchentliche Reviews, in denen wir die Aufgaben nach Bedarf neu verteilt haben.

Das Commit-Bild bestätigt das: beide durchgängig aktiv über die gesamte Laufzeit. Durch die wöchentliche Neuverteilung hat jeder auch im Code des anderen gearbeitet — das machte Reviews leichter und verhindert, dass Wissen bei einer Person hängen bleibt.

Ich selbst habe den ersten Prototyp aufgesetzt und früh die Richtung „maximal smart plus sehr modernes UI" getrieben — agentischer Chat, Chat-Gedächtnis und das Aurora-Glass-Design, das Sie übrigens auch in diesen Folien sehen.

Ehrliche Einordnung zum Schluss: Bei zwei Personen ist dieses Modell empfehlenswert. In größeren Teams würden wir es kritisch hinterfragen — dort braucht es klarere Verantwortlichkeiten.""")

# ---------------------------------------------------------------- 6 · Q4
s = new_slide("bg_content.png")
top = q_header(s, "04", "Welche Tools haben Sie genutzt?",
               "Mehrere KI-Assistenten im Vergleich getestet — Claude Code (Opus) klar am stärksten; dazu ein schlanker Produkt-Stack.")
yend = bullet_stack(s, 220, 1480, [
    ([W("GitHub Copilot (Student) — "), N("Limits schnell aufgebraucht; Qualität nur für kleine Aufgaben überzeugend.")], False),
    ([W("Codex & Gemini (Free) — "), N("gute bis brauchbare Ergebnisse, Free-Tier jeweils schnell erschöpft.")], False),
    ([W("Claude Code (Opus, Paid) — "), N("am stärksten: ganze Umbauten in einem Rutsch. Aber: lange Laufzeiten (>10 Min) & abrupte Token-Stops.")], True),
], top=top + 20, bottom=760, cap=40)
eyebrow(s, 220, yend + 34, "Produkt-Stack")
chip(s, 220, yend + 70, 1480,
     [{"t": "Groq + OpenRouter · Supabase · Vitest · GitHub Actions · MkDocs Material · Flask", "f": INM, "c": FG}])
hud(s, 5)
NOTES.append("""⏱ ~2 Min · Sprecher: Ujjwal

Ich habe die Coding-Assistenten selbst gegeneinander getestet — im Editor war durchgehend VS Code gesetzt.

GitHub Copilot im Studenten-Tarif: Die Limits waren trotzdem schnell aufgebraucht, und die Qualität hat bei größeren Aufgaben nicht überzeugt — für Kleinkram okay.

Codex und Gemini im Free-Tier lieferten brauchbare bis gute Ergebnisse — Gemini schneller und etwas besser als Copilot —, aber die Kontingente waren rasch erschöpft.

Klar am stärksten war Claude Code mit Opus: Es hat ganze Umbauten in einem einzigen Durchlauf erledigt — zum Beispiel das komplette Aurora-Glass-Redesign unserer Projektwebseite inklusive neuem SVG-Logo. Der Preis: Es arbeitet oft länger als zehn Minuten, und Token-Limits stoppen den Prozess trotz Paid-Version abrupt — weiterarbeiten geht erst nach dem Reset.

Im Produkt selbst setzen wir bewusst auf einen schlanken Stack: Groq und OpenRouter als LLM-APIs, Supabase als Datenbank, Vitest für Tests, GitHub Actions für CI, MkDocs für die Webseite, Flask für den Doku-Agenten.""")

# ---------------------------------------------------------------- 7 · Q5
s = new_slide("bg_content.png")
top = q_header(s, "05", "Wo haben Sie die meiste Zeit verloren?",
               "Token- und Rate-Limits, eine anfangs unklare Produktvision — und die Zuverlässigkeit des Form-Fillings über fremde Seiten.")
bullet_stack(s, 220, 1480, [
    ([W("Entwicklungspausen durch Token- & Rate-Limits "), N("(v. a. Claude Code).")], False),
    ([W("Unklare Produktvision: "), N("teils in gegensätzliche Richtungen entwickelt — Teile mussten neu gebaut werden.")], False),
    ([W("Anforderungen strukturieren "), N("+ Projektmanagement, Doku und alles „drumherum“.")], False),
    ([W("Zuverlässigkeit des Form-Fillings: "), N("mehrere Agent-Iterationen, bis Profil & gelernte Daten verlässlich genutzt wurden.")], False),
    ([W("War Story: "), N("GitHub-Push-Protection blockierte einen hardcodierten API-Key — gute Lektion zu Secrets-Management.")], True),
], top=top, gap=12)
hud(s, 6)
NOTES.append("""⏱ ~2 Min · Sprecher: Maximilian (optionale Ergänzung: Ujjwal)

Am meisten Zeit verloren haben wir an vier Stellen.

Erstens: Entwicklungspausen durch Token- und Rate-Limits — vor allem bei Claude Code, wo ein Limit den Prozess mitten in der Arbeit stoppt.

Zweitens — und das ist selbstkritisch —: eine anfangs unklare Produktvision. Wir waren uns über Richtung und Ausbau nicht immer einig, haben teils in gegensätzliche Richtungen entwickelt und mussten Teile neu bauen.

Drittens: das Strukturieren der Anforderungen und alles „drumherum“ — Projektmanagement, Dokumentation.

Viertens: die Zuverlässigkeit des Form-Fillings über fremde Seiten — es brauchte mehrere Agent-Iterationen, bis Profil und gelernte Daten verlässlich genutzt wurden.

Und eine War Story zum Schmunzeln: GitHubs Push-Protection hat einen versehentlich hardcodierten API-Key blockiert. Der Umbau kostete Zeit, war aber eine gute Lektion in Secrets-Management.

[Optional, Ujjwal:] Ausgezahlt hat sich dagegen der Refactor vom großen content.js in sieben Module — danach war jede Änderung lokal und testbar. Das hätten wir rückblickend früher gemacht.""")

# ---------------------------------------------------------------- 8 · Q6
s = new_slide("bg_content.png")
top = q_header(s, "06", "Was hat Ihnen für bessere Ergebnisse gefehlt?",
               "Eine früher klar definierte Vorgehensweise — und mehr konkrete Beispiele und Referenz-Architekturen in der Vorlesung.")
bullet_stack(s, 220, 1480, [
    ([W("Eine von Anfang an klar definierte Vorgehensweise "), N("hätte uns Richtungswechsel und doppelte Arbeit erspart.")], False),
    ([W("„Wo Sie nach der Einheit stehen sollten“ war hilfreich — "), N("noch hilfreicher wäre eine klarere Richtung gewesen.")], False),
    ([W("Konkrete Referenz-Architekturen "), N("(z. B. „Browser-Extension + LLM“) als frühe Orientierung.")], True),
], top=top, gap=16)
hud(s, 7)
NOTES.append("""⏱ ~1 Min · Sprecher: Ujjwal

Ehrliche Antwort: vor allem eine von Anfang an klar definierte Vorgehensweise — die hätte uns einen Teil der Richtungswechsel erspart, die Maximilian gerade bei Frage 5 beschrieben hat.

Aus der Vorlesung war das Format „Wo Sie nach der Einheit stehen sollten“ sehr hilfreich — das war unser Kompass für den Projektstand.

Noch hilfreicher wären eine klarere inhaltliche Richtung und mehr konkrete Beispiele gewesen — etwa Referenz-Architekturen wie „Browser-Extension plus LLM-Backend“, an denen man sich früh orientieren kann, statt die Grundarchitektur selbst mehrfach umzubauen.""")

# ---------------------------------------------------------------- 9 · Q7
s = new_slide("bg_content.png")
top = q_header(s, "07", "Ihre größten Lerneffekte?",
               "Klassische Methodik wird durch KI wichtiger, nicht überflüssig — und „KI managen“ ist eine eigene Kompetenz.")
bullet_stack(s, 220, 1480, [
    ([W("Anforderungen, Stakeholder, Mockups — "), N("durch KI "), W("noch wichtiger"), N(" geworden.")], False),
    ([W("Nachvollziehbarkeit & Doku der KI-Arbeit "), N("ist eigener Aufwand → Memory-/Entscheidungslog + autonomer Doku-Agent.")], False),
    ([W("„KI managen“ ist eine eigene Kompetenz: "), N("Kontext geben, Ergebnisse reviewen, Entscheidungen festhalten.")], False),
    ([W("Kontext-Investment zahlt sich aus: "), N("Regeldatei (CLAUDE.md), Personas & Memory-Log → spürbar bessere, konsistente Ergebnisse.")], True),
    ([W("Modernes UI entsteht iterativ — "), N("Aurora-Glass wurde über mehrere Runden geschärft, nicht im ersten Wurf.")], False),
], top=top, gap=12)
hud(s, 8)
NOTES.append("""⏱ ~1,5 Min · Sprecher: Ujjwal

Unser wichtigster Lerneffekt klingt zunächst paradox: Klassische Methodik wird durch KI wichtiger, nicht überflüssig. Anforderungen definieren, Stakeholder einbinden, Mockups bauen — all das braucht es erst recht, wenn die KI in Minuten Code produziert. Sonst produziert sie schnell das Falsche.

Zweitens: Die Nachvollziehbarkeit der KI-Arbeit ist ein eigener Aufwand. Wir haben dafür ein Memory- und Entscheidungslog geführt — und schließlich den autonomen Doku-Agenten gebaut, der das nach jedem Commit selbst erledigt.

Drittens: „KI managen“ ist eine eigene Kompetenz — Kontext geben, Ergebnisse reviewen, Entscheidungen festhalten.

Mein persönlicher Aha-Moment: Kontext-Investment schlägt bessere Prompts. Erst mit der Regeldatei CLAUDE.md, den Personas und dem Memory-Log wurden die KI-Ergebnisse wirklich konsistent — vorher hat die KI in jeder Sitzung Entscheidungen neu erfunden.

Und ganz praktisch: Ein modernes UI entsteht iterativ. Das Aurora-Glass-Design haben wir über mehrere Runden geschärft — nicht im ersten Wurf.""")

# ---------------------------------------------------------------- 10 · Q8
s = new_slide("bg_content.png")
top = q_header(s, "08", "Ihre Empfehlung an eine Firma?",
               "Unbedingt nutzen — aber Architektur, Review und Guardrails bleiben menschlich; Provider-Kosten von Tag 1 einplanen.")
bullet_stack(s, 220, 1480, [
    ([W("In Stunden zum Prototyp — "), N("Prototyping "), W("ohne"), N(" KI ist nicht mehr zu empfehlen.")], False),
    ([W("Ergebnisqualität pro Zeitaufwand sehr hoch — "), N("der Prototyp taugt als Basis für die Weiterentwicklung.")], False),
    ([W("Architektur- & Richtungsentscheidungen sowie Review bleiben menschlich.")], True),
    ([W("Guardrails & Sicherheit gehören schon in den Prototyp: "), N("keine Keys im Code, kein automatisches Absenden.")], False),
    ([W("Provider-Kosten & Limits von Tag 1 einplanen: "), N("Free-Tiers sind schnell leer, auch Paid limitiert — Fallbacks & Budget vorsehen.")], False),
], top=top, gap=12)
hud(s, 9)
NOTES.append("""⏱ ~1,5 Min · Sprecher: Maximilian

Unsere Empfehlung an eine Firma: unbedingt nutzen.

Man kommt in wenigen Stunden zu einem funktionierenden Prototyp — Prototyping ohne KI ist aus unserer Sicht schlicht nicht mehr zu empfehlen. Die Ergebnisqualität im Verhältnis zum Zeitaufwand ist sehr hoch, und der Prototyp taugt als echte Basis für die Weiterentwicklung.

Aber drei Dinge gehören dazu:

Erstens: Architektur- und Richtungsentscheidungen sowie das Review bleiben menschlich. Die KI liefert Code — die Verantwortung liefert sie nicht mit.

Zweitens: Guardrails und Sicherheit gehören schon in den Prototyp, nicht erst in die Produktion — bei uns: keine API-Keys im Code und kein automatisches Absenden.

Drittens: Provider-Kosten und Limits von Tag 1 einplanen. Free-Tiers sind schnell aufgebraucht, auch Paid-Tiers limitieren — also Fallbacks und Budget vorsehen.

[Optional, Ujjwal:] Aus eigener Erfahrung: Ein automatischer Provider-Fallback gehört für mich in jeden ernsthaften Prototyp — nicht erst in die Produktion.""")

# ---------------------------------------------------------------- 11 · Bilanz
s = new_slide("bg_content.png")
eyebrow(s, 220, 118, "Grenzen · Stolz · Ausblick")
text(s, 220, 152, 1200, 60, [{"t": "Ehrliche Bilanz", "f": SGS, "c": "FFFFFF"}], 26)
text(s, 220, 262, 710, 40, [{"t": "Worauf wir stolz sind", "f": SGS, "c": PINKISH}], 15)
y = 316
for segs, hi in [
    ([W("Batch-Agent: "), N("1 statt 12 API-Calls pro Formularseite.")], False),
    ([W("Auto-Fallback Groq → OpenRouter, "), N("transparent per Toast.")], False),
    ([W("~77 % Coverage + CI "), N("bei Vanilla JS ohne Build-Step.")], False),
    ([W("Autonomer Doku-Agent "), N("dokumentiert via Git-Hook nach jedem Commit.")], False),
]:
    y = chip(s, 220, y, 710, segs, size=11.5, hi=hi) + 16
text(s, 990, 262, 710, 40, [{"t": "Grenzen & Ausblick", "f": SGS, "c": PINKISH}], 15)
y = 316
for segs, hi in [
    ([W("Client-seitiger API-Key = Prototyp-Niveau "), N("→ produktiv: Backend-Proxy + Consent-Flow.")], False),
    ([W("Nicht erreichbar: "), N("cross-origin iFrames & Chrome-PDF-Viewer.")], False),
    ([W("Ausblick: "), N("Vision-OCR · RAG über persönliche Dokumente (pgvector) · MCP-Server.")], True),
]:
    y = chip(s, 990, y, 710, segs, size=11.5, hi=hi) + 16
link_pill(s, 560, 820, 800, 64)
hud(s, 10)
NOTES.append("""⏱ ~1 Min · Sprecher: beide (Reserve für Rückfragen)

Zum Abschluss eine ehrliche Bilanz.

Worauf wir stolz sind: der Batch-Agent — ein API-Call statt zwölf pro Formularseite; der automatische Provider-Fallback von Groq zu OpenRouter, transparent per Toast; rund 77 Prozent Test-Coverage mit CI — bei einer Vanilla-JS-Extension ohne Build-Step; und der autonome Doku-Agent, der sich per Git-Hook selbst um die Dokumentation kümmert.

Genauso klar benennen wir die Grenzen: Der client-seitige API-Key ist Prototyp-Niveau — produktiv bräuchte es einen Backend-Proxy und einen Consent-Flow pro Formular. Cross-origin iFrames und der native Chrome-PDF-Viewer sind technisch nicht erreichbar.

Der Ausblick: Vision-OCR — also Ausweis fotografieren statt abtippen —, echtes RAG über persönliche Dokumente mit pgvector, und ein MCP-Server für den Doku-Agenten.

Alles davon ist öffentlich dokumentiert — auf unserer Stakeholder-Webseite.""")

# ---------------------------------------------------------------- 12 · Danke
s = new_slide("bg_title.png")
s.shapes.add_picture(f"{ASSETS}/logo.png", PX(960 - 80), PX(248), PX(160), PX(160))
text(s, 260, 448, 1400, 160, [{"t": "Danke — Fragen?", "f": SG, "b": True, "grad": TITLE_GRAD}],
     52, align=PP_ALIGN.CENTER)
link_pill(s, 560, 650, 800, 66, size=14.5)
text(s, 360, 768, 1200, 36,
     [{"t": "FormAssist · Maximilian Plitzko & Ujjwal Makkar · AI-Prototyping SS 2026", "f": IN, "c": DIM}],
     11.5, align=PP_ALIGN.CENTER)
hud(s, 11)
NOTES.append("""⏱ ~30 Sek · Sprecher: beide

Vielen Dank für Ihre Aufmerksamkeit!

Alles, was wir heute gezeigt haben, ist öffentlich dokumentiert — Architektur, Entscheidungen, Teststrategie und Projektstand — auf unserer Stakeholder-Webseite unter ujjwalmak.github.io/form-ai-assistant.

Wir freuen uns auf Ihre Fragen.

(Tipp: Bei Rückfragen zu Grenzen, Selbstkritik oder Ausblick → zurück zu Folie 11. Bei Technik-Detailfragen: Batch-Agent, Fallback-Logik und Test-Setup sind in memory/decisions.md und auf der Webseite dokumentiert.)""")

# ---- Notizen anhängen ----
for slide, note in zip(prs.slides, NOTES):
    slide.notes_slide.notes_text_frame.text = note

prs.save(OUT)
print("saved:", OUT)
