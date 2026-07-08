#!/usr/bin/env python3
"""FormAssist — Abschluss: reine Tool-Demo (Aurora-Glass), 10 Folien, 16:9.
Drehbuch + Backup für die Live-Demo am 09.07.2026 — mit echten Prototyp-Screenshots
(assets/demo/, aufgenommen via Playwright; KI-Antworten für die Aufnahme gemockt)."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls

import os
_HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(_HERE, "assets")
OUT = os.path.join(_HERE, "..", "FormAssist_Abschluss.pptx")

EMU_PX = 6350  # 12192000 / 1920
def PX(v): return Emu(int(v * EMU_PX))

# ---- Design-Tokens (aus reflexion.html) ----
GRAD = [(0, "8B5CF6"), (48, "D946EF"), (100, "F472B6")]
TITLE_GRAD = [(10, "FFFFFF"), (42, "E9D5FF"), (72, "F0ABFC"), (100, "F9A8D4")]
FG, DIM, LILAC, PINKISH = "ECEBF2", "B0AAC2", "E9D5FF", "F0ABFC"
SG, SGM, SGS = "Space Grotesk", "Space Grotesk Medium", "Space Grotesk SemiBold"
IN, INM, INS = "Inter", "Inter Medium", "Inter SemiBold"
ANG_110, ANG_V = 1200000, 5400000

FILLS = ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill")

def _gradfill_xml(stops, ang):
    gs = "".join(f'<a:gs pos="{int(p*1000)}"><a:srgbClr val="{c}"/></a:gs>' for p, c in stops)
    return (f'<a:gradFill {nsdecls("a")} rotWithShape="1"><a:gsLst>{gs}</a:gsLst>'
            f'<a:lin ang="{ang}" scaled="1"/></a:gradFill>')

def _set_sp_fill(shape, xml):
    spPr = shape._element.spPr
    for tag in FILLS:
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    ln = spPr.find(qn("a:ln"))
    el = parse_xml(xml)
    if ln is not None:
        ln.addprevious(el)
    else:
        spPr.append(el)

def solid_alpha(shape, hexc, alpha_pct):
    _set_sp_fill(shape, f'<a:solidFill {nsdecls("a")}><a:srgbClr val="{hexc}">'
                        f'<a:alpha val="{int(alpha_pct*1000)}"/></a:srgbClr></a:solidFill>')

def grad_fill(shape, stops=GRAD, ang=ANG_110):
    _set_sp_fill(shape, _gradfill_xml(stops, ang))

def _get_ln(shape):
    spPr = shape._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = parse_xml(f'<a:ln {nsdecls("a")}/>')
        spPr.append(ln)
    for tag in FILLS:
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    return ln

def line_alpha(shape, hexc, alpha_pct, w_pt):
    ln = _get_ln(shape)
    ln.set("w", str(int(w_pt * 12700)))
    ln.insert(0, parse_xml(f'<a:solidFill {nsdecls("a")}><a:srgbClr val="{hexc}">'
                           f'<a:alpha val="{int(alpha_pct*1000)}"/></a:srgbClr></a:solidFill>'))

def grad_line(shape, stops=GRAD, ang=ANG_110, w_pt=1.2):
    ln = _get_ln(shape)
    ln.set("w", str(int(w_pt * 12700)))
    ln.insert(0, parse_xml(_gradfill_xml(stops, ang)))

def no_line(shape):
    _get_ln(shape).append(parse_xml(f'<a:noFill {nsdecls("a")}/>'))

def glow(shape, hexc="D946EF", alpha_pct=55, rad_pt=6):
    shape._element.spPr.append(parse_xml(
        f'<a:effectLst {nsdecls("a")}><a:glow rad="{int(rad_pt*12700)}">'
        f'<a:srgbClr val="{hexc}"><a:alpha val="{int(alpha_pct*1000)}"/></a:srgbClr>'
        f'</a:glow></a:effectLst>'))

def no_shadow(shape):
    shape.shadow.inherit = False

def grad_run(run, stops=GRAD, ang=ANG_110):
    rPr = run._r.get_or_add_rPr()
    for tag in FILLS:
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    el = parse_xml(_gradfill_xml(stops, ang))
    ln = rPr.find(qn("a:ln"))
    if ln is not None:
        ln.addnext(el)
    else:
        rPr.insert(0, el)

def add_runs(p, segs, size, default_font=IN, default_color=FG):
    for s in segs:
        r = p.add_run()
        r.text = s["t"]
        r.font.size = Pt(s.get("sz", size))
        r.font.name = s.get("f", default_font)
        r.font.bold = s.get("b", False)
        if s.get("spc"):
            r._r.get_or_add_rPr().set("spc", str(int(s["spc"])))
        if s.get("grad"):
            grad_run(r, s["grad"] if isinstance(s["grad"], list) else GRAD)
        else:
            r.font.color.rgb = RGBColor.from_string(s.get("c", default_color))
        if s.get("href"):
            r.hyperlink.address = s["href"]
    return p

def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(PX(x), PX(y), PX(w), PX(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb, tf

def text(slide, x, y, w, h, segs, size, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=None,
         default_font=IN, default_color=FG):
    tb, tf = textbox(slide, x, y, w, h, anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    if ls:
        p.line_spacing = ls
    add_runs(p, segs, size, default_font, default_color)
    return tb

def est_lines(segs, size, inner_w):
    n = sum(len(s["t"]) for s in segs)
    cpl = inner_w / (size * 2 * 0.48)
    return 1 if n <= cpl else 2

def chip_h(segs, size, w):
    lines = est_lines(segs, size, w - 56 - 26)
    return (58 if size >= 12 else 54) if lines == 1 else (88 if size >= 12 else 84)

def chip(slide, x, y, w, segs, size=12.5, hi=False, dot=True, ls=1.3, align=PP_ALIGN.LEFT):
    h = chip_h(segs, size, w)
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sh.adjustments[0] = min(0.5, 14.0 / h)
    solid_alpha(sh, "FFFFFF", 6.0 if hi else 3.5)
    if hi:
        grad_line(sh, w_pt=1.4)
    else:
        line_alpha(sh, "FFFFFF", 7, 1.0)
    no_shadow(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left, tf.margin_right = PX(56 if dot else 26), PX(26)
    tf.margin_top = tf.margin_bottom = PX(6)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = ls
    add_runs(p, segs, size)
    if dot:
        d = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x + 24), PX(y + h/2 - 5.5), PX(11), PX(11))
        d.adjustments[0] = 0.33
        grad_fill(d)
        no_line(d)
        glow(d)
        no_shadow(d)
    return y + h

def bullet_stack(slide, x, w, specs, top, bottom=945, gap=14, bias=0.40, cap=120):
    """specs: list of (segs, hi) — vertikal ausbalanciert im Kartenbereich."""
    total = sum(chip_h(segs, 12.5, w) for segs, _ in specs) + gap * (len(specs) - 1)
    y = top + max(0, min(cap, (bottom - top - total) * bias))
    for segs, hi in specs:
        y = chip(slide, x, y, w, segs, hi=hi) + gap
    return y - gap

def eyebrow(slide, x, y, txt, color=LILAC, size=10, align=PP_ALIGN.LEFT, w=900):
    return text(slide, x, y, w, 30, [{"t": txt.upper(), "f": INS, "c": color, "spc": 110}],
                size, align=align)

def stat_tile(slide, x, y, w, h, num, label):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sh.adjustments[0] = min(0.5, 18.0 / h)
    solid_alpha(sh, "FFFFFF", 3.5)
    line_alpha(sh, "FFFFFF", 8, 1.0)
    no_shadow(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = PX(12)
    tf.margin_top = tf.margin_bottom = PX(8)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    add_runs(p1, [{"t": num, "f": SG, "b": True, "grad": True}], 30)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    add_runs(p2, [{"t": label, "f": IN, "c": DIM}], 10.5)

def mini_card(slide, x, y, w, h, title, sub, hi=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sh.adjustments[0] = min(0.5, 16.0 / h)
    solid_alpha(sh, "FFFFFF", 6.0 if hi else 3.5)
    if hi:
        grad_line(sh, w_pt=1.4)
    else:
        line_alpha(sh, "FFFFFF", 8, 1.0)
    no_shadow(sh)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = PX(10)
    tf.margin_top = tf.margin_bottom = PX(6)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    add_runs(p1, [{"t": title, "f": INS, "c": "FFFFFF"}], 13)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(3)
    add_runs(p2, [{"t": sub, "f": IN, "c": DIM}], 9.5)

def flow_row(slide, y, cards, x=220, w=331, h=112, gap=52):
    """cards: list of (title, sub, hi) — horizontale Kette mit Gradient-Pfeilen."""
    for i, (t1, t2, hi) in enumerate(cards):
        mini_card(slide, x + i * (w + gap), y, w, h, t1, t2, hi=hi)
        if i < len(cards) - 1:
            text(slide, x + i * (w + gap) + w, y, gap, h,
                 [{"t": "→", "f": SG, "b": True, "grad": True}], 20,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

TOTAL = 10

def hud(slide, idx, total=TOTAL):
    if idx > 0:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PX(0), PX(0), PX(1920 * idx / (total - 1)), PX(6))
        grad_fill(bar)
        no_line(bar)
        no_shadow(bar)
    text(slide, 1650, 1032, 230, 30,
         [{"t": f"{idx+1:02d}", "f": INS, "c": "FFFFFF"}, {"t": f" / {total:02d}", "f": IN, "c": DIM}],
         9.5, align=PP_ALIGN.RIGHT)

def header(slide, brow, title, key=None):
    eyebrow(slide, 220, 118, brow)
    text(slide, 220, 152, 1480, 60, [{"t": title, "f": SGS, "c": "FFFFFF"}], 26)
    if key is None:
        return 252
    lines = est_lines([{"t": key}], 15.5, 1420)
    bar_h = 46 if lines == 1 else 72
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(221), PX(232), PX(5), PX(bar_h))
    bar.adjustments[0] = 0.5
    grad_fill(bar, ang=ANG_V)
    no_line(bar)
    no_shadow(bar)
    text(slide, 246, 234, 1434, 80, [{"t": key, "f": INS, "c": "F6F3FB"}], 15.5, ls=1.32)
    return 352 if lines == 1 else 366

REPO_URL = "https://github.com/ujjwalmak/form-ai-assistant"
SITE_URL = "https://ujjwalmak.github.io/form-ai-assistant/"

def link_pill(slide, x, y, w, h, label, href, size=14):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    pill.adjustments[0] = 0.5
    grad_fill(pill)
    no_line(pill)
    no_shadow(pill)
    # Klick-Ziel auf der Form statt Text-Hyperlink — so bleibt der Text weiß
    # (Run-Hyperlinks bekämen PowerPoints blaues Unterstrich-Styling).
    pill.click_action.hyperlink.address = href
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = PX(10)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_runs(p, [{"t": label, "f": INS, "c": "FFFFFF"}], size)

W = lambda t: {"t": t, "f": INS, "c": "FFFFFF"}   # weißes Semibold-Segment
N = lambda t: {"t": t}                              # normales Segment

# =====================================================================
prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
prs.core_properties.title = "FormAssist — Tool-Demo · AI-Prototyping SS 2026"
prs.core_properties.author = "Maximilian Plitzko · Ujjwal Makkar"
BLANK = prs.slide_layouts[6]
NOTES = []
DEMO = os.path.join(ASSETS, "demo")

def new_slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(f"{ASSETS}/{bg}", 0, 0, Emu(12192000), Emu(6858000))
    return s

def demo_slide(idx_hud, kap, total_kap, title, img, points, note):
    """Ein Demo-Kapitel: großer Screenshot links, „Live zeigen“-Chips rechts."""
    s = new_slide("bg_content.png")
    eyebrow(s, 220, 118, f"Live-Demo · Kapitel {kap}/{total_kap}")
    text(s, 220, 152, 1480, 60, [{"t": title, "f": SGS, "c": "FFFFFF"}], 26)
    # gerahmter Screenshot (Canvas 3400x2000 inkl. Glow-Rand)
    s.shapes.add_picture(f"{DEMO}/{img}", PX(160), PX(268), PX(1096), PX(645))
    eyebrow(s, 1310, 300, "Live zeigen", size=9.5, w=440)
    y = 344
    for segs, hi in points:
        y = chip(s, 1310, y, 440, segs, size=11, hi=hi) + 14
    hud(s, idx_hud)
    NOTES.append(note)
    return s

# ---------------------------------------------------------------- 1 · Titel
s = new_slide("bg_title.png")
s.shapes.add_picture(f"{ASSETS}/logo.png", PX(960 - 92), PX(148), PX(184), PX(184))
badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(720), PX(352), PX(480), PX(46))
badge.adjustments[0] = 0.5
solid_alpha(badge, "D946EF", 12)
line_alpha(badge, "D946EF", 40, 1.0)
no_shadow(badge)
btf = badge.text_frame
btf.vertical_anchor = MSO_ANCHOR.MIDDLE
btf.margin_left = btf.margin_right = PX(10)
btf.margin_top = btf.margin_bottom = 0
bp = btf.paragraphs[0]
bp.alignment = PP_ALIGN.CENTER
add_runs(bp, [{"t": "TOOL-DEMO · AI-PROTOTYPING SS 2026", "f": INS, "c": LILAC, "spc": 90}], 10.5)
text(s, 160, 418, 1600, 200, [{"t": "FormAssist", "f": SG, "b": True, "grad": TITLE_GRAD}],
     66, align=PP_ALIGN.CENTER)
text(s, 460, 634, 1000, 80,
     [{"t": "Der KI-Agent, der komplexe Browser-Formulare versteht, erklärt und ausfüllt — datensparsam und nie ohne Freigabe.",
       "f": IN, "c": DIM}], 14.5, align=PP_ALIGN.CENTER, ls=1.4)
text(s, 360, 758, 1200, 40,
     [{"t": "Team  ", "f": INS, "c": "FFFFFF"},
      {"t": "Maximilian Plitzko · Ujjwal Makkar", "f": INS, "grad": True}],
     14, align=PP_ALIGN.CENTER)
text(s, 360, 806, 1200, 34,
     [{"t": "Prof. Dr. Sebastian Dünnebeil · FK07 · 09.07.2026", "f": IN, "c": DIM}],
     11.5, align=PP_ALIGN.CENTER)
hud(s, 0)
NOTES.append("""⏱ ~30 Sek · Sprecher: Maximilian (Einstieg)

Guten Tag zusammen! Wir sind Maximilian Plitzko und Ujjwal Makkar — und das ist FormAssist: unser KI-Agent, der komplexe Browser-Formulare versteht, in einfacher Sprache erklärt und ausfüllt.

Heute gibt es bewusst wenig Folien: Das hier ist eine reine Tool-Demo. Wir zeigen den Prototyp die nächsten Minuten live im Browser — in sieben kurzen Kapiteln.""")

# ---------------------------------------------------------------- 2 · Fahrplan
s = new_slide("bg_content.png")
top = header(s, "Der Fahrplan", "7 Kapitel — alles live im Browser",
             "Die Folien sind nur das Drehbuch (und der Plan B): Jedes Kapitel zeigt echte Screenshots des Prototyps.")
row1 = [("01 · Erkennen", "Sidebar & Feld-Badge", False),
        ("02 · Erklären", "in einfacher Sprache", False),
        ("03 · Profil", "alle Daten an einem Ort", False),
        ("04 · ⚡ Agent", "füllt & fragt gezielt", False)]
row2 = [("05 · Multi-Page", "navigiert selbst weiter", False),
        ("06 · Chat", "Befehle direkt ausführen", False),
        ("07 · Submit-Review", "der Mensch sendet ab", True)]
fw, fh, gap = 331, 112, 52
for i, (t1, t2, hi) in enumerate(row1):
    mini_card(s, 220 + i * (fw + gap), 430, fw, fh, t1, t2, hi=hi)
x2 = (1920 - (3 * fw + 2 * gap)) / 2
for i, (t1, t2, hi) in enumerate(row2):
    mini_card(s, x2 + i * (fw + gap), 582, fw, fh, t1, t2, hi=hi)
text(s, 460, 780, 1000, 34,
     [{"t": "Test-Umgebung: eigene Test-Site (6 Formulartypen) · Extension als entpackte Erweiterung geladen", "f": IN, "c": DIM}],
     11, align=PP_ALIGN.CENTER)
hud(s, 1)
NOTES.append("""⏱ ~30 Sek · Sprecher: Maximilian

Der Fahrplan für die Demo — sieben Kapitel: Die Sidebar erkennt das Formular, erklärt es, unser Profil liefert die Daten, der Agent füllt aus und fragt gezielt nach, navigiert selbst über mehrere Seiten, der Chat führt Befehle direkt aus — und am Ende der wichtigste Punkt: Abgesendet wird nur vom Menschen.

Wir zeigen das auf unserer eigenen Test-Site mit sechs Formulartypen — die Extension läuft als ganz normale entpackte Chrome-Erweiterung.

[Übergabe an Ujjwal, Browser-Screenshare starten. Ab hier: Ujjwal fährt, Maximilian moderiert.]

>>> AB HIER LIVE IM BROWSER — die folgenden Folien sind das Backup, falls die Live-Demo klemmt. <<<""")

# ---------------------------------------------------------------- 3–9 · Kapitel
demo_slide(2, 1, 7, "Die Sidebar kennt das Formular", "demo1_erkennen.png", [
    ([W("Trigger-Button mit Feldanzahl"), N(" — Alt+Shift+F öffnet die Sidebar.")], False),
    ([W("7 Felder erkannt: "), N("Label, Typ, Pflichtfeld — semantisch gelesen, nicht nur HTML.")], False),
    ([W("Trust-Row: "), N("Profil lokal gespeichert · KI nur bei Aktionen.")], False),
], """⏱ ~1 Min · Fahrer: Ujjwal · Moderation: Maximilian

LIVE: Registration-Wizard auf der Test-Site öffnen. Trigger-Button unten rechts zeigen (Badge = erkannte Felder). Sidebar öffnen (Klick oder Alt+Shift+F). Feldliste aufklappen: Label, Typ, Pflicht-Badges.

SAGEN (Maximilian): FormAssist scannt das Formular semantisch — Labels, Hinweise, Fehlermeldungen, auch im Shadow DOM. Unten die Trust-Row: Das Profil bleibt lokal, die KI wird nur bei Aktionen gefragt.

(Hinweis zu den Folien-Screenshots: echte Extension-UI auf der Test-Site, automatisiert aufgenommen; die KI-Antworten wurden für die Aufnahme deterministisch gemockt — live antwortet Groq.)""")

demo_slide(3, 2, 7, "Erklären statt Behördendeutsch", "demo2_erklaeren.png", [
    ([W("Ein Klick auf „Erklären“ — "), N("die KI liest das komplette Formular.")], False),
    ([W("Zweck · Pflichtfelder · Stolperstellen"), N(" — kompakt als Chat-Antwort.")], False),
    ([W("Live-Kontext: "), N("vor jeder Antwort wird die Seite neu gescannt.")], False),
], """⏱ ~45 Sek · Fahrer: Ujjwal · Moderation: Maximilian

LIVE: „Erklären“ klicken. Antwort kurz vorlesen lassen: Zweck, Pflichtfelder, Stolperstellen.

SAGEN: Das ist der Anti-Behördendeutsch-Knopf — das Formular in einfacher Sprache. Vor jeder Antwort scannt FormAssist die Seite neu, die KI sieht also immer den aktuellen Zustand.""")

demo_slide(4, 3, 7, "Das Profil: alle Daten an einem Ort", "demo3_profil.png", [
    ([W("15 Standardfelder"), N(" (Person, Adresse, Kontakt, Bank, Beruf) + gelernte Extras.")], False),
    ([W("Mehrere Profile, "), N("Import/Export als JSON — lokal in chrome.storage.")], False),
    ([W("Optional: "), N("Supabase-Sync geräteübergreifend.")], False),
], """⏱ ~45 Sek · Fahrer: Ujjwal · Moderation: Maximilian

LIVE: Profil-Panel öffnen (Kopfzeile). Durch die Felder scrollen: 15 Standardfelder, komplett gefüllt. Kurz auf Extras & Profil-Switcher zeigen.

SAGEN: Hier ist die Antwort auf „die Daten liegen verstreut": einmal pflegen, überall ausfüllen. Alles lokal in chrome.storage — optional per Supabase über Geräte synchronisiert.""")

demo_slide(5, 4, 7, "⚡ Der Agent füllt — und fragt gezielt", "demo4_agent.png", [
    ([W("Direkt aus Profil: 3 Felder"), N(" — sofort, ganz ohne API-Call.")], False),
    ([W("KI nur für Unbekanntes, "), N("gebatcht: 1 Call für bis zu 12 Felder.")], True),
    ([W("Wirklich Unklares → "), N("gezielte Rückfrage als Chip — keine Raterei.")], False),
], """⏱ ~1,5 Min · Fahrer: Ujjwal · Moderation: Maximilian

LIVE: Zurück zum Formular, ⚡ „Formular ausfüllen“ klicken. Zeigen, wie Name + Geburtsdatum sofort dastehen („Direkt aus Profil: 3 Felder“), Gender per KI gewählt wird — und der Agent bei „Middle Name“ gezielt nachfragt statt zu raten.

SAGEN (Maximilian): Drei Stufen: Erst deterministisch aus dem Profil — null API-Kosten. Dann die KI, aber gebatcht: ein einziger Call für bis zu zwölf unbekannte Felder. Und was wirklich niemand wissen kann, wird gefragt — als Chip, präzise pro Feld. Nach dem Füllen prüft der Agent übrigens die Validierungsfehler der Seite und korrigiert selbst nach.""")

demo_slide(6, 5, 7, "Multi-Page: der Agent navigiert selbst", "demo5_multipage.png", [
    ([W("Antwort tippen — "), N("der Agent füllt und macht selbst weiter.")], False),
    ([W("Klickt „Next“, füllt Seite 2 "), N("— Schritt für Schritt durch den Wizard.")], False),
    ([W("Merkt sich Antworten "), N("für Folgeseiten · Loop-Schutz nach 12 Seiten.")], False),
], """⏱ ~1 Min · Fahrer: Ujjwal · Moderation: Maximilian

LIVE: Die Rückfrage beantworten (z. B. „Alexander“ tippen). Der Agent bedankt sich, klickt selbst auf „Next" und füllt Schritt 2 — E-Mail und Telefon stehen schon da, der Fortschrittsbalken wandert.

SAGEN: Das ist der Unterschied zwischen Autofill und einem Agenten: Er arbeitet sich selbstständig durch mehrseitige Formulare, merkt sich Ihre Antworten für Folgeseiten — und hat einen Loop-Schutz nach zwölf Seiten.""")

demo_slide(7, 6, 7, "Der Chat kann zupacken", "demo6_chat.png", [
    ([W("„Trag bei E-Mail … ein“ — "), N("die KI führt die Aktion direkt im Formular aus.")], False),
    ([W("Bestätigung mit Feld + Wert"), N(" — jede Aktion nachvollziehbar.")], False),
    ([W("Gedächtnis pro Domain: "), N("der Chat erinnert sich über Seitenwechsel & Sessions.")], False),
], """⏱ ~1 Min · Fahrer: Ujjwal · Moderation: Maximilian

LIVE: Auf der Checkout-Seite in den Chat tippen: „Trag bei E-Mail demo@formassist.de ein“. Zeigen: Die KI antwortet UND das Feld ist gefüllt — mit Bestätigungs-Karte (Feld + Wert).

SAGEN: Der Chat ist nicht nur Auskunft, er kann zupacken — fill, select, check, Radios per Optionstext. Nur eine Aktion kennt er nicht: submit. Dazu gleich mehr. Und: Das Chat-Gedächtnis pro Domain überlebt Seitenwechsel und Sessions.""")

demo_slide(8, 7, 7, "Submit-Review: der Mensch sendet ab", "demo7_review.png", [
    ([W("Absenden wird abgefangen: "), N("Status OK · Warnung · Fehlt.")], False),
    ([W("„Trotzdem absenden“ — "), N("der letzte Klick gehört immer dem Menschen.")], True),
    ([W("Harte Guardrail: "), N("der Action-Parser kennt kein „submit“.")], False),
], """⏱ ~1 Min · Fahrer: Ujjwal · Moderation: Maximilian

LIVE: Formular absenden klicken → FormAssist fängt ab und prüft: Status OK, kurze Begründung, dazu „Erneut prüfen" / „Trotzdem absenden". Bewusst auf „Trotzdem absenden" klicken.

SAGEN (Maximilian, Punchline): Und das ist der wichtigste Klick der ganzen Demo — den macht immer der Mensch. Es gibt kein automatisches Absenden; der Action-Parser kennt die Aktion schlicht nicht. Guardrails sind bei uns Architektur, kein Feature.

[Zurück zu den Folien für den Abschluss.]""")

# ---------------------------------------------------------------- 10 · Danke + Links
s = new_slide("bg_title.png")
s.shapes.add_picture(f"{ASSETS}/logo.png", PX(960 - 80), PX(200), PX(160), PX(160))
text(s, 260, 396, 1400, 140, [{"t": "Danke — probieren Sie es aus!", "f": SG, "b": True, "grad": TITLE_GRAD}],
     46, align=PP_ALIGN.CENTER)
link_pill(s, 560, 574, 800, 64, "⭐  github.com/ujjwalmak/form-ai-assistant", REPO_URL, size=13.5)
link_pill(s, 560, 662, 800, 64, "🌐  ujjwalmak.github.io/form-ai-assistant", SITE_URL, size=13.5)
text(s, 360, 790, 1200, 36,
     [{"t": "FormAssist · Maximilian Plitzko & Ujjwal Makkar · AI-Prototyping SS 2026", "f": IN, "c": DIM}],
     11.5, align=PP_ALIGN.CENTER)
hud(s, 9)
NOTES.append("""⏱ ~30 Sek · Sprecher: beide

Vielen Dank! Repo und Projektwebseite sind hier verlinkt — der Repo-Link kommt jetzt in den Zoom-Chat, dann können Sie FormAssist direkt selbst laden: Repo klonen, chrome://extensions, „Entpackte Erweiterung laden“, API-Key eintragen — fertig. Auf der Webseite stehen Architektur, Entscheidungen und Teststrategie.

Wir freuen uns auf Ihre Fragen.

(Orga: Foliensatz bis So 12.07. auf Moodle hochladen — Links sind auf dieser Folie. Falls das Repo privat ist: „duenne“ auf GitHub einladen oder Repo auf Moodle hochladen.
Rückfragen-Reserve: Technik → Kapitel 4/5 · Sicherheit & Grenzen → Kapitel 7 · Tests: 69 Unit-Tests, ~77 % Coverage, CI · Doku-Agent & MkDocs-Webseite → siehe Projektwebseite.)""")

# ---- Notizen anhängen ----
for slide, note in zip(prs.slides, NOTES):
    slide.notes_slide.notes_text_frame.text = note

prs.save(OUT)
print("saved:", OUT)
